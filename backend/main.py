from fastapi import FastAPI, UploadFile, Form, HTTPException
from fastapi.responses import JSONResponse, FileResponse
from fastapi.middleware.cors import CORSMiddleware
import httpx
import pdfplumber
import docx
from pptx import Presentation
import uuid
import os

app = FastAPI()

# Allow CORS if needed
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# n8n webhook URL
N8N_WEBHOOK_URL = "http://localhost:5678/webhook-test/upload"
TEMP_DIR = "temp_uploads"
os.makedirs(TEMP_DIR, exist_ok=True)

async def extract_text(file: UploadFile) -> str:
    ext = file.filename.split(".")[-1].lower()
    temp_path = f"{TEMP_DIR}/{uuid.uuid4()}.{ext}"
    
    try:
        # Save temporarily for processing
        with open(temp_path, "wb") as buffer:
            buffer.write(await file.read())
        
        if ext == "pdf":
            text = ""
            with pdfplumber.open(temp_path) as pdf:
                for page in pdf.pages:
                    text += page.extract_text() or ""
            return text
        
        elif ext == "docx":
            document = docx.Document(temp_path)
            return "\n".join([p.text for p in document.paragraphs])
        
        elif ext == "pptx":
            presentation = Presentation(temp_path)
            text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        
        else:
            with open(temp_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
    
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

@app.post("/api/upload")
async def upload_assignment(
    file: UploadFile,
    assignment_type: str = Form(...),
    user_email: str = Form(...)
):
    try:
        extracted_text = await extract_text(file)
        
        async with httpx.AsyncClient() as client:
            response = await client.post(
                N8N_WEBHOOK_URL,
                json={
                    "content": extracted_text,
                    "assignmentType": assignment_type,
                    "userEmail": user_email,
                    "fileName": file.filename
                },
                timeout=30
            )
        
        if response.status_code != 200:
            raise HTTPException(status_code=500, detail="Processing failed")
            
        return JSONResponse({
            "status": "success",
            "message": "Assignment submitted for processing"
        })
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))